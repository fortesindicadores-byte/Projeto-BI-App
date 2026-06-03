from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import re

# ─── Paleta ───────────────────────────────────────────────
BG      = RGBColor(0x0C, 0x10, 0x17)
CARD    = RGBColor(0x14, 0x1B, 0x26)
CARD2   = RGBColor(0x1A, 0x23, 0x35)
BORDER  = RGBColor(0x1E, 0x2D, 0x40)
ORANGE  = RGBColor(0xF9, 0x73, 0x16)
BLUE    = RGBColor(0x38, 0xBD, 0xF8)
TEXT    = RGBColor(0xF1, 0xF5, 0xF9)
TEXT2   = RGBColor(0x94, 0xA3, 0xB8)
TEXT3   = RGBColor(0x47, 0x55, 0x69)
GREEN   = RGBColor(0x3B, 0xB3, 0x3B)
RED     = RGBColor(0xFF, 0x66, 0x66)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

W = Inches(13.33)  # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank = prs.slide_layouts[6]  # blank layout

def add_slide():
    return prs.slides.add_slide(blank)

def bg(slide, color=BG):
    bg_shape = slide.shapes.add_shape(1, 0, 0, W, H)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = color
    bg_shape.line.fill.background()

def txt(slide, text, x, y, w, h, size=14, bold=False, color=TEXT,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    run.font.name = 'Montserrat'
    return txBox

def rect(slide, x, y, w, h, fill=CARD, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def orange_bar(slide, x=0, y=0, w=None, h=0.04):
    w = w or 13.33
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ORANGE
    bar.line.fill.background()

def pill(slide, text, x, y, w=1.6, h=0.3, fill=CARD2, tc=TEXT2, size=9):
    r = rect(slide, x, y, w, h, fill=fill, line_color=BORDER, line_width=Pt(0.75))
    txt(slide, text, x+0.08, y+0.03, w-0.1, h-0.05, size=size, color=tc, align=PP_ALIGN.CENTER)

def section_label(slide, label, x=0.45, y=0.22):
    txt(slide, label, x, y, 5, 0.3, size=9, bold=True, color=ORANGE)
    orange_bar(slide, x, y+0.28, 2.5, 0.025)

# ══════════════════════════════════════════════════════════
# SLIDE 1 — CAPA
# ══════════════════════════════════════════════════════════
s = add_slide()
bg(s)
# Gradient rectangles (simulated)
grad = s.shapes.add_shape(1, Inches(0), Inches(0), Inches(4), H)
grad.fill.solid(); grad.fill.fore_color.rgb = RGBColor(0x12, 0x18, 0x22)
grad.line.fill.background()

# Logo mark
r = rect(s, 0.55, 1.8, 0.85, 0.85, fill=RGBColor(0x2A, 0x1A, 0x0C), line_color=ORANGE, line_width=Pt(1))
txt(s, "BI", 0.55, 1.88, 0.85, 0.7, size=22, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

txt(s, "GESTÃO EM MOVIMENTO", 1.55, 1.82, 8, 0.5, size=26, bold=True, color=TEXT)
txt(s, "Plataforma de Business Intelligence — Fortes Indicadores", 1.55, 2.32, 9, 0.4, size=12, color=TEXT2)
orange_bar(s, 1.55, 2.78, 8, 0.04)
txt(s, "Referência Técnica Completa · 2026", 1.55, 2.9, 8, 0.35, size=10, color=TEXT3)

# Stat pills
pills = [
    ("5 painéis ativos", 1.55, 3.6),
    ("Google Sheets + Supabase", 3.35, 3.6),
    ("Chart.js 4.4", 5.75, 3.6),
    ("HTML puro — sem framework", 7.3, 3.6),
]
for label, px, py in pills:
    pill(s, label, px, py, w=1.6, tc=TEXT2)

txt(s, "fortesindicadores-byte.github.io/gestao-em-movimento", 1.55, 6.8, 10, 0.4, size=9, color=TEXT3, italic=True)

# ══════════════════════════════════════════════════════════
# SLIDE 2 — VISÃO GERAL DA ARQUITETURA
# ══════════════════════════════════════════════════════════
s = add_slide()
bg(s)
orange_bar(s)
section_label(s, "ARQUITETURA", 0.45, 0.2)
txt(s, "Como o sistema funciona", 0.45, 0.55, 10, 0.35, size=18, bold=True, color=TEXT)

# 3 colunas: Dados → Painéis → Usuários
cols = [
    ("DADOS\n(Google Sheets)", "Planilha compartilhada\nEstrutura: nv3, mês, valor\nAcesso: GVIZ API (público)\nSem backend necessário", 0.5, ORANGE),
    ("PAINÉIS\n(GitHub Pages)", "HTML + CSS + JS puro\nChart.js 4.4 + Datalabels\nAutocontido (1 arquivo)\nGitHub Pages — CDN global", 4.5, BLUE),
    ("ACESSO\n(Supabase Auth)", "Login / Cadastro\nAprovação manual\nToken refresh ~55 min\nFix: flag isApproved", 8.5, GREEN),
]
for title, body, cx, color in cols:
    r = rect(s, cx, 1.2, 4.0, 4.0, fill=CARD, line_color=BORDER, line_width=Pt(0.75))
    # colored top bar
    tb = s.shapes.add_shape(1, Inches(cx), Inches(1.2), Inches(4.0), Inches(0.07))
    tb.fill.solid(); tb.fill.fore_color.rgb = color; tb.line.fill.background()
    txt(s, title, cx+0.15, 1.35, 3.7, 0.6, size=13, bold=True, color=color)
    txt(s, body, cx+0.15, 2.05, 3.7, 2.8, size=10, color=TEXT2)

# Arrow connectors (visual only — rectangles)
for ax in [4.42, 8.42]:
    ar = s.shapes.add_shape(1, Inches(ax), Inches(3.0), Inches(0.06), Inches(0.06))
    ar.fill.solid(); ar.fill.fore_color.rgb = TEXT3; ar.line.fill.background()
txt(s, "→", 4.35, 2.95, 0.5, 0.4, size=18, color=TEXT3, align=PP_ALIGN.CENTER)
txt(s, "→", 8.35, 2.95, 0.5, 0.4, size=18, color=TEXT3, align=PP_ALIGN.CENTER)

txt(s, "⚠  git push origin main → 503 (proxy antigo)  |  Sempre usar MCP para publicar no main", 0.45, 5.55, 12.4, 0.4, size=9, color=RED, bold=True)

# ══════════════════════════════════════════════════════════
# SLIDE 3 — PAINÉIS ATIVOS
# ══════════════════════════════════════════════════════════
s = add_slide()
bg(s)
orange_bar(s)
section_label(s, "PAINÉIS", 0.45, 0.2)
txt(s, "Painéis ativos e roadmap", 0.45, 0.55, 10, 0.35, size=18, bold=True, color=TEXT)

panels = [
    ("Visão Financeira", "DRE Consolidado por conta × mês\nFiltros: Ano, Mês, Conta, Projeto, Unidade\nGráficos: linha (Real/Rem/Orc), barra mensal\nTabela: DRE full com totais", ORANGE, "ATIVO"),
    ("Painel KM", "KM realizado × remunerado × orçado\nFiltros: Ano, Mês, Projeto\nGráficos: linha, Dispersão KM %, Δ% barras\nHero: 52px, cards 5 colunas", ORANGE, "ATIVO"),
    ("R$/KM", "Custo por km por projeto\nFiltros: Ano, Mês, Projeto\nTabela: impacto descrescente por desvio\nGráfico: Δ% horizontal + linha tendência", ORANGE, "ATIVO"),
    ("Árvore Combustível", "Consumo por projeto/tipo combustível\nFiltros: Ano, Mês, Projeto, Tipo\nGráfico: barras agrupadas (litros + R$)\nNested: /combustivel/arvore-combustivel/", GREEN, "ATIVO"),
    ("Financeiro Pessoal", "Controle Renan & Tati\nAcesso direto (não aparece no hub)\nURL: /financeiro-pessoal/", BLUE, "ATIVO"),
    ("Eficiência Km/L\nPreço R$/L\nConsumo CO²\n+ 7 outros", "Pastas vazias (.gitkeep)\nA criar conforme demanda\nSeguir padrão visao-financeira", TEXT3, "EM BREVE"),
]

for i, (name, desc, color, status) in enumerate(panels):
    col = i % 3
    row = i // 3
    cx = 0.45 + col * 4.3
    cy = 1.1 + row * 2.8
    r = rect(s, cx, cy, 4.0, 2.5, fill=CARD, line_color=BORDER, line_width=Pt(0.75))
    tb = s.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(4.0), Inches(0.06))
    tb.fill.solid(); tb.fill.fore_color.rgb = color; tb.line.fill.background()
    # status badge
    sc = GREEN if status == "ATIVO" else TEXT3
    txt(s, f"● {status}", cx+0.12, cy+0.1, 3.8, 0.25, size=8, color=sc, bold=True)
    txt(s, name, cx+0.12, cy+0.38, 3.8, 0.45, size=11, bold=True, color=color)
    txt(s, desc, cx+0.12, cy+0.85, 3.8, 1.5, size=9, color=TEXT2)

# ══════════════════════════════════════════════════════════
# SLIDE 4 — PALETA + TIPOGRAFIA
# ══════════════════════════════════════════════════════════
s = add_slide()
bg(s)
orange_bar(s)
section_label(s, "DESIGN SYSTEM", 0.45, 0.2)
txt(s, "Paleta de cores e tipografia", 0.45, 0.55, 10, 0.35, size=18, bold=True, color=TEXT)

colors_left = [
    ("--bg",     "#0C1017", BG,     TEXT2),
    ("--card",   "#141B26", CARD,   TEXT2),
    ("--card2",  "#1A2335", CARD2,  TEXT2),
    ("--border", "#1E2D40", BORDER, TEXT2),
    ("--orange", "#F97316", ORANGE, WHITE),
    ("--blue",   "#38BDF8", BLUE,   BG),
    ("--text",   "#F1F5F9", TEXT,   BG),
    ("--text2",  "#94A3B8", TEXT2,  BG),
    ("--text3",  "#475569", TEXT3,  WHITE),
    ("--green",  "#3BB33B", GREEN,  WHITE),
    ("--red",    "#FF6666", RED,    WHITE),
]
for i, (varname, hex_v, fill, tc) in enumerate(colors_left):
    col = i % 2
    row = i // 2
    cx = 0.45 + col * 3.1
    cy = 1.15 + row * 0.56
    r = rect(s, cx, cy, 2.85, 0.48, fill=fill, line_color=BORDER, line_width=Pt(0.5))
    txt(s, varname, cx+0.1, cy+0.05, 1.6, 0.38, size=9, bold=True, color=tc)
    txt(s, hex_v, cx+1.7, cy+0.09, 1.0, 0.3, size=8, color=tc, align=PP_ALIGN.RIGHT)

# Typography table
txt(s, "TIPOGRAFIA", 6.9, 1.1, 6, 0.3, size=9, bold=True, color=ORANGE)
orange_bar(s, 6.9, 1.38, 6.0, 0.025)

typo = [
    ("Hero value",          "52px", "800"),
    ("Título seção/tabela", "16px", "800"),
    ("Título de card",      "15px", "800"),
    ("Card value",          "24px", "800"),
    ("Card label",          "10px", "600"),
    ("Dados da tabela",     "12-13px", "400-500"),
    ("Cabeçalho coluna",    "11px", "700"),
    ("Row total",           "12-13px", "700"),
    ("Badges / hints",      "9-10px", "400"),
]
for i, (elem, sz, wt) in enumerate(typo):
    cy = 1.5 + i * 0.5
    r = rect(s, 6.9, cy, 6.0, 0.46, fill=CARD if i % 2 == 0 else CARD2)
    txt(s, elem, 7.0, cy+0.1, 3.0, 0.3, size=9, color=TEXT2)
    txt(s, sz,   10.1, cy+0.1, 1.2, 0.3, size=9, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    txt(s, wt,   11.4, cy+0.1, 1.3, 0.3, size=9, color=TEXT2, align=PP_ALIGN.CENTER)

txt(s, "Fonte: Montserrat (Google Fonts) — pesos 400/500/600/700/800", 6.9, 6.1, 6, 0.3, size=8, color=TEXT3, italic=True)

# ══════════════════════════════════════════════════════════
# SLIDE 5 — FONTE DE DADOS + PADRÃO nv3
# ══════════════════════════════════════════════════════════
s = add_slide()
bg(s)
orange_bar(s)
section_label(s, "DADOS", 0.45, 0.2)
txt(s, "Google Sheets como fonte de dados", 0.45, 0.55, 10, 0.35, size=18, bold=True, color=TEXT)

# GVIZ box
r = rect(s, 0.45, 1.1, 6.0, 2.5, fill=CARD, line_color=BORDER, line_width=Pt(0.75))
txt(s, "GVIZ API — Endpoint público", 0.6, 1.15, 5.7, 0.35, size=11, bold=True, color=ORANGE)
code = ("const SHEET_ID  = '<id-planilha>';\n"
        "const SHEET_TAB = '<aba>';\n"
        "const URL = `https://docs.google.com/spreadsheets/d/\n"
        "  ${SHEET_ID}/gviz/tq?tqx=out:json&sheet=${encodeURIComponent(SHEET_TAB)}`;\n\n"
        "// Busca:\n"
        "const text = await (await fetch(URL)).text();\n"
        "const json = JSON.parse(text.match(/setResponse\\(([\\s\\S]*)\\)/)[1]);\n"
        "const rows = json.table.rows;  // rows[i].c[j].v = valor")
txt(s, code, 0.6, 1.55, 5.7, 2.0, size=8, color=TEXT2)

# nv3 pattern box
r = rect(s, 0.45, 3.75, 6.0, 2.8, fill=CARD, line_color=BORDER, line_width=Pt(0.75))
txt(s, "Padrão de campo nv3", 0.6, 3.82, 5.7, 0.3, size=11, bold=True, color=ORANGE)
txt(s, '"PROJETO - UNIDADE (INATIVO)"', 0.6, 4.2, 5.7, 0.35, size=12, bold=True, color=BLUE)
examples = [
    ('• "LOGÍSTICA - BLC (INATIVO)"', '"BLC" (strip INATIVO)'),
    ('• "TRANSPORTE - SP1"',          '"SP1"'),
    ('• "FROTA - CAB"',               '"CAB"'),
]
for i, (raw, parsed) in enumerate(examples):
    cy = 4.65 + i * 0.4
    txt(s, raw,    0.65, cy, 3.5, 0.35, size=9, color=TEXT2)
    txt(s, "→ label: " + parsed, 4.1, cy, 2.3, 0.35, size=9, color=GREEN)
txt(s, "getNv3Prefix(v) → tudo antes de '-'  |  getUniLabel(v) → tudo após '-', sem (INATIVO)", 0.6, 6.2, 5.7, 0.3, size=8, color=TEXT3, italic=True)

# Right side: paneis e abas
r = rect(s, 6.7, 1.1, 6.2, 5.45, fill=CARD, line_color=BORDER, line_width=Pt(0.75))
txt(s, "Planilhas por painel", 6.85, 1.18, 5.9, 0.3, size=11, bold=True, color=ORANGE)

planilhas = [
    ("Visão Financeira",    "DRE Consolidado"),
    ("Painel KM",           "KM"),
    ("R$/KM",               "RS por KM"),
    ("Árvore Combustível",  "Combustivel"),
    ("Financeiro Pessoal",  "aba própria"),
]
for i, (painel, aba) in enumerate(planilhas):
    cy = 1.6 + i * 0.6
    rb = rect(s, 6.7, cy, 6.2, 0.55, fill=CARD2 if i % 2 else CARD)
    txt(s, painel, 6.85, cy+0.1, 3.2, 0.35, size=9, bold=True, color=TEXT)
    pill(s, aba, 10.2, cy+0.1, 2.5, 0.32, fill=RGBColor(0x0a, 0x0f, 0x18), tc=BLUE, size=8)

txt(s, "⚠  Planilha deve ser pública (Compartilhar → Qualquer pessoa com o link → Leitor)", 6.85, 4.9, 5.9, 0.5, size=9, color=RED)
txt(s, "rows[i].c[j].v = valor bruto  |  rows[i].c[j].f = valor formatado", 6.85, 5.5, 5.9, 0.35, size=8, color=TEXT3, italic=True)

# ══════════════════════════════════════════════════════════
# SLIDE 6 — FILTROS MULTI-SELECT
# ══════════════════════════════════════════════════════════
s = add_slide()
bg(s)
orange_bar(s)
section_label(s, "FILTROS", 0.45, 0.2)
txt(s, "Sistema de filtros multi-select", 0.45, 0.55, 10, 0.35, size=18, bold=True, color=TEXT)

# Left: visual filter mockup
r = rect(s, 0.45, 1.1, 5.5, 5.5, fill=CARD, line_color=BORDER, line_width=Pt(0.75))
txt(s, "Estrutura de filtro", 0.6, 1.18, 5.2, 0.3, size=11, bold=True, color=ORANGE)

# Button mockup
rb = rect(s, 0.65, 1.6, 2.0, 0.42, fill=RGBColor(0x0a, 0x0f, 0x18), line_color=RGBColor(0x2a, 0x3a, 0x50), line_width=Pt(1))
txt(s, "ANO  [2]", 0.72, 1.67, 1.85, 0.3, size=9, bold=True, color=TEXT)
# badge
bb = rect(s, 2.4, 1.66, 0.3, 0.28, fill=ORANGE)
txt(s, "2", 2.4, 1.67, 0.3, 0.28, size=8, bold=True, color=BG, align=PP_ALIGN.CENTER)

# Dropdown mockup
drop = rect(s, 0.65, 2.05, 2.5, 1.8, fill=RGBColor(0x0f, 0x18, 0x24), line_color=RGBColor(0x2a, 0x3a, 0x50), line_width=Pt(0.75))
opts = [("☑ Todos", True), ("☑ 2025", True), ("☑ 2024", True), ("☐ 2023", False)]
for i, (opt, chk) in enumerate(opts):
    cy = 2.1 + i * 0.4
    bg_opt = CARD2 if i == 0 else RGBColor(0x0f, 0x18, 0x24)
    ob = rect(s, 0.65, cy, 2.5, 0.38, fill=bg_opt)
    tc = TEXT if chk else TEXT3
    txt(s, opt, 0.77, cy+0.06, 2.3, 0.28, size=9, color=tc, bold=(i==0))

txt(s, "• Checkbox 'Todos' sincroniza todos\n• Badge mostra qtd selecionada\n• Clique fora fecha dropdown\n• Múltiplos filtros simultâneos", 0.65, 4.0, 4.8, 1.5, size=9, color=TEXT2)

# Right: filter table
r = rect(s, 6.2, 1.1, 6.7, 5.5, fill=CARD, line_color=BORDER, line_width=Pt(0.75))
txt(s, "Filtros por painel", 6.35, 1.18, 6.4, 0.3, size=11, bold=True, color=ORANGE)

filter_data = [
    ("Visão Financeira",   ["Ano", "Mês", "Conta", "Projeto", "Unidade"],   5),
    ("Painel KM",          ["Ano", "Mês", "Projeto"],                        3),
    ("R$/KM",              ["Ano", "Mês", "Projeto"],                        3),
    ("Árvore Combustível", ["Ano", "Mês", "Projeto", "Tipo Combustível"],    4),
]
for i, (painel, filtros, n) in enumerate(filter_data):
    cy = 1.65 + i * 1.1
    txt(s, painel, 6.35, cy, 6.4, 0.3, size=10, bold=True, color=TEXT)
    for j, f in enumerate(filtros):
        px = 6.35 + j * 1.25
        pill(s, f, px, cy+0.35, w=1.15, h=0.28, fill=RGBColor(0x0a, 0x0f, 0x18), tc=BLUE, size=8)
    txt(s, f"Mobile: {'3' if n >= 5 else '2'} por linha", 6.35, cy+0.72, 6.4, 0.25, size=8, color=TEXT3, italic=True)

# ══════════════════════════════════════════════════════════
# SLIDE 7 — TIPOS DE GRÁFICOS
# ══════════════════════════════════════════════════════════
s = add_slide()
bg(s)
orange_bar(s)
section_label(s, "GRÁFICOS", 0.45, 0.2)
txt(s, "Tipos de gráficos — Chart.js 4.4", 0.45, 0.55, 10, 0.35, size=18, bold=True, color=TEXT)

charts = [
    ("Linha — Evolução Mensal",
     "Datasets: Realizado (laranja sólido, bw=3)\nRemunerado (azul tracejado [5,3], bw=2)\nOrçado (branco/cinza tracejado, bw=2)\nPontos destacados: mês selecionado = radius 6, branco",
     ORANGE),
    ("Barras Verticais — Dispersão %",
     "Orientação: indexAxis padrão (vertical)\nCores: laranja (positivo) / vermelho (negativo)\nDatalabels: ±X.X% com 1 decimal\nUsado em: Painel KM — Dispersão de KM %",
     BLUE),
    ("Barras Horizontais — Δ% Projetos",
     "indexAxis: 'y'\nCores: laranja (positivo) / vermelho (negativo)\nDatalabels: anchor end/start conforme sinal\nFormatter: Math.round (inteiro) — sem decimal",
     GREEN),
    ("Barras Agrupadas — Combustível",
     "Grupos por tipo (diesel/arla/gasolina)\nDatalabels duplos: litros + R$\nNested: /combustivel/arvore-combustivel/",
     ORANGE),
    ("Cards Sparkline",
     "Mini linha dentro de card KPI\nSem eixos, sem legend, sem tooltip\ntension: 0.3, pointRadius: 0\nApenas tendência visual",
     BLUE),
]

for i, (title, desc, color) in enumerate(charts):
    col = i % 2 if i < 4 else 0
    row = i // 2
    cx = 0.45 + col * 6.4
    cy = 1.1 + row * 2.0
    if i == 4:  # last centered
        cx = 0.45; cy = 5.2
        w = 12.4
    else:
        w = 5.9
    r = rect(s, cx, cy, w, 1.85, fill=CARD, line_color=BORDER, line_width=Pt(0.75))
    tb = s.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(w), Inches(0.06))
    tb.fill.solid(); tb.fill.fore_color.rgb = color; tb.line.fill.background()
    txt(s, title, cx+0.12, cy+0.12, w-0.2, 0.3, size=10, bold=True, color=color)
    txt(s, desc, cx+0.12, cy+0.5, w-0.2, 1.3, size=9, color=TEXT2)

# ══════════════════════════════════════════════════════════
# SLIDE 8 — TIPOS DE TABELAS
# ══════════════════════════════════════════════════════════
s = add_slide()
bg(s)
orange_bar(s)
section_label(s, "TABELAS", 0.45, 0.2)
txt(s, "Tipos e padrões de tabelas", 0.45, 0.55, 10, 0.35, size=18, bold=True, color=TEXT)

# DRE table mockup (visao-financeira)
r = rect(s, 0.45, 1.1, 6.2, 4.2, fill=CARD, line_color=BORDER, line_width=Pt(0.75))
txt(s, "DRE Consolidado (Visão Financeira)", 0.6, 1.17, 5.9, 0.3, size=10, bold=True, color=ORANGE)
# Header row
hr = rect(s, 0.45, 1.55, 6.2, 0.38, fill=RGBColor(0x0a, 0x0f, 0x18))
headers = ["CONTA", "JAN", "FEV", "MAR", "...", "TOTAL"]
for j, h in enumerate(headers):
    hx = 0.6 + j * 0.95
    txt(s, h, hx, 1.6, 0.9, 0.25, size=7, bold=True, color=TEXT3, align=PP_ALIGN.CENTER)
# Data rows
rows_data = [
    ("Receita Bruta",   ["850k", "920k", "880k", "...", "10.2 mi"], GREEN),
    ("Deduções",        ["-45k", "-52k", "-48k", "...", "-580k"],   RED),
    ("Receita Líquida", ["805k", "868k", "832k", "...", "9.6 mi"],  TEXT),
    ("Custos",          ["-620k","-680k","-650k","...", "-7.8 mi"], RED),
    ("EBITDA",          ["185k", "188k", "182k", "...", "1.8 mi"],  ORANGE),
]
for i, (conta, vals, color) in enumerate(rows_data):
    cy = 2.0 + i * 0.48
    rb = rect(s, 0.45, cy, 6.2, 0.46, fill=CARD2 if i % 2 else CARD)
    txt(s, conta, 0.6, cy+0.1, 1.5, 0.28, size=9, color=TEXT if color == TEXT else color)
    for j, v in enumerate(vals):
        txt(s, v, 1.7 + j * 0.95, cy+0.1, 0.9, 0.28, size=9, color=color, align=PP_ALIGN.RIGHT)
# Total
tr2 = rect(s, 0.45, 4.45, 6.2, 0.45, fill=RGBColor(0x1E, 0x2D, 0x40))
txt(s, "TOTAL  →  13.6 mi", 0.6, 4.52, 5.9, 0.32, size=9, bold=True, color=ORANGE)
txt(s, "• Fonte: 12–13px  • Hover: rgba(255,255,255,.035)\n• Total: border-top 2px, font-weight 700\n• Colunas numéricas: text-align right", 0.6, 5.0, 5.9, 0.7, size=8, color=TEXT3)

# Impact table (rs-por-km)
r = rect(s, 6.9, 1.1, 6.1, 4.2, fill=CARD, line_color=BORDER, line_width=Pt(0.75))
txt(s, "Impacto R$/KM (R$/KM)", 7.05, 1.17, 5.8, 0.3, size=10, bold=True, color=ORANGE)
hr2 = rect(s, 6.9, 1.55, 6.1, 0.38, fill=RGBColor(0x0a, 0x0f, 0x18))
h2 = ["PROJETO", "REAL", "REM", "Δ", "Δ%", "IMPACTO"]
for j, h in enumerate(h2):
    hx = 7.05 + j * 0.9
    txt(s, h, hx, 1.6, 0.85, 0.25, size=7, bold=True, color=TEXT3, align=PP_ALIGN.CENTER)
impact_rows = [
    ("LOGÍSTICA",   "2.45", "2.20", "+0.25", "+11.4%", "+18.2k", RED),
    ("TRANSPORTE",  "3.10", "2.95", "+0.15", "+5.1%",  "+9.8k",  RED),
    ("FROTA LEVE",  "1.85", "1.90", "-0.05", "-2.6%",  "-3.1k",  GREEN),
    ("AGREGADOS",   "4.20", "4.35", "-0.15", "-3.4%",  "-5.9k",  GREEN),
]
for i, (proj, real, rem, d, dp, imp, color) in enumerate(impact_rows):
    cy = 2.0 + i * 0.48
    rb = rect(s, 6.9, cy, 6.1, 0.46, fill=CARD2 if i % 2 else CARD)
    vals = [proj, real, rem, d, dp, imp]
    colors = [TEXT, TEXT2, TEXT2, color, color, color]
    for j, (v, vc) in enumerate(zip(vals, colors)):
        txt(s, v, 7.05 + j * 0.9, cy+0.1, 0.85, 0.28, size=9, color=vc, align=PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT)
txt(s, "• Ordenado por |impacto| decrescente\n• .imp-cell: font-size 13px\n• Vermelho = custo > remunerado (ruim)", 7.05, 5.0, 5.8, 0.7, size=8, color=TEXT3)

# ══════════════════════════════════════════════════════════
# SLIDE 9 — AUTENTICAÇÃO SUPABASE
# ══════════════════════════════════════════════════════════
s = add_slide()
bg(s)
orange_bar(s)
section_label(s, "AUTH", 0.45, 0.2)
txt(s, "Autenticação Supabase + Hub", 0.45, 0.55, 10, 0.35, size=18, bold=True, color=TEXT)

# Flow diagram
flows = [
    ("Login\nsignInWithPassword", ORANGE),
    ("Cadastro\nsignUp → e-mail confirm", BLUE),
    ("Esqueci Senha\nresetPasswordForEmail", TEXT2),
    ("Redefinir\nupdateUser({password})\nevento PASSWORD_RECOVERY", GREEN),
]
for i, (flow, color) in enumerate(flows):
    cx = 0.45 + i * 3.1
    r = rect(s, cx, 1.1, 2.85, 1.4, fill=CARD, line_color=color, line_width=Pt(1))
    txt(s, flow, cx+0.12, 1.2, 2.65, 1.2, size=9, bold=True, color=color)

# Key fix box
r = rect(s, 0.45, 2.75, 12.4, 1.6, fill=RGBColor(0x1a, 0x10, 0x08), line_color=ORANGE, line_width=Pt(1))
txt(s, "🔧 FIX CRÍTICO: Crash do hub — Supabase dispara SIGNED_IN a cada ~55 min (refresh de token)", 0.6, 2.82, 12.0, 0.3, size=10, bold=True, color=ORANGE)
fix_code = ("let isApproved = false;           // flag global\n"
            "function showHub(user) { isApproved = true; ... }      // setar true ao entrar\n"
            "function showAuth()    { isApproved = false; ... }     // setar false ao sair\n"
            "// onAuthStateChange:\n"
            "if (event==='SIGNED_IN' && session && !isRegistering && !isApproved) await checkApproval(session.user);\n"
            "// catch em checkApproval:\n"
            "if (!isApproved) showAuth();      // só expulsa se não estava no hub")
txt(s, fix_code, 0.6, 3.15, 12.0, 1.1, size=8, color=TEXT2)

# Hub screens
r = rect(s, 0.45, 4.55, 5.8, 2.6, fill=CARD, line_color=BORDER, line_width=Pt(0.75))
txt(s, "Telas do Hub", 0.6, 4.62, 5.5, 0.3, size=11, bold=True, color=ORANGE)
screens = [
    ("Auth Screen (#auth-screen)", "Login / Cadastro / Esqueci / Redefinir"),
    ("Pending Screen",              "Aguardando aprovação na tabela approved_users"),
    ("Hub Screen (#hub-screen)",    "Grid de cards por cluster (após aprovação)"),
]
for i, (name, desc) in enumerate(screens):
    cy = 5.05 + i * 0.65
    txt(s, f"• {name}", 0.62, cy, 5.5, 0.28, size=9, bold=True, color=BLUE)
    txt(s, desc, 0.62, cy+0.28, 5.5, 0.3, size=8, color=TEXT2)

# Supabase config
r = rect(s, 6.5, 4.55, 6.4, 2.6, fill=CARD, line_color=BORDER, line_width=Pt(0.75))
txt(s, "Config Supabase Dashboard", 6.65, 4.62, 6.1, 0.3, size=11, bold=True, color=ORANGE)
cfg = [
    ("Site URL",    "https://fortesindicadores-byte.github.io/gestao-em-movimento/"),
    ("SMTP",        "Resend — smtp.resend.com:465, user=resend"),
    ("E-mail conf", "Ativado (obrigatório confirmar para logar)"),
    ("Tabela DB",   "approved_users — email column — acesso manual"),
]
for i, (k, v) in enumerate(cfg):
    cy = 5.05 + i * 0.55
    txt(s, k + ":", 6.65, cy, 1.5, 0.3, size=9, bold=True, color=TEXT2)
    txt(s, v, 8.2, cy, 4.6, 0.3, size=8, color=TEXT)

# ══════════════════════════════════════════════════════════
# SLIDE 10 — WORKFLOW DE PUSH
# ══════════════════════════════════════════════════════════
s = add_slide()
bg(s)
orange_bar(s)
section_label(s, "GIT / DEPLOY", 0.45, 0.2)
txt(s, "Como publicar mudanças — Push Workflow", 0.45, 0.55, 10, 0.35, size=18, bold=True, color=TEXT)

# Warning
r = rect(s, 0.45, 1.0, 12.4, 0.55, fill=RGBColor(0x2a, 0x0c, 0x0c), line_color=RED, line_width=Pt(1))
txt(s, "⚠  git push origin main → 503  |  Repositório renomeado (Projeto-BI-App → gestao-em-movimento). Proxy local desatualizado.", 0.6, 1.1, 12.0, 0.35, size=9, bold=True, color=RED)

# Steps
steps = [
    ("1", "Editar arquivo local", "Edit tool no Claude Code\nEx: visao-financeira/index.html", ORANGE),
    ("2", "Commit no dev branch", "git add <arquivo>\ngit commit -m 'descrição'\ngit push -u origin claude/great-allen-OVObS", BLUE),
    ("3", "Buscar SHA do arquivo", "mcp__github__get_file_contents\n  owner: fortesindicadores-byte\n  repo: gestao-em-movimento\n  path: <caminho>\n→ retorna .sha atual", GREEN),
    ("4", "Push para main via MCP", "mcp__github__create_or_update_file\n  owner, repo, path, message\n  content: <conteúdo do arquivo>\n  sha: <sha do passo 3>", ORANGE),
    ("5", "Sincronizar local", "git fetch origin main\ngit checkout origin/main -- <arquivo>\ngit add && git commit && git push dev", BLUE),
]
for i, (num, title, code, color) in enumerate(steps):
    col = i % 3 if i < 3 else (i - 3)
    row = 0 if i < 3 else 1
    cx = 0.45 + col * 4.3
    cy = 1.75 + row * 2.7
    if i >= 3:
        cx = 0.45 + (i - 3) * 6.3
    r = rect(s, cx, cy, 4.1 if i < 3 else 5.8, 2.5, fill=CARD, line_color=color, line_width=Pt(1))
    # number bubble
    nb = rect(s, cx+0.12, cy+0.12, 0.35, 0.35, fill=color)
    txt(s, num, cx+0.12, cy+0.12, 0.35, 0.35, size=10, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(s, title, cx+0.6, cy+0.15, 3.3, 0.3, size=10, bold=True, color=color)
    txt(s, code, cx+0.15, cy+0.55, 3.8, 1.85, size=8, color=TEXT2)

# MCP Auth note
r = rect(s, 0.45, 7.0, 12.4, 0.38, fill=CARD2)
txt(s, "Se MCP pedir autenticação: mcp__github__authenticate → URL → usuário autoriza → copiar URL callback → mcp__github__complete_authentication(callback_url)", 0.6, 7.07, 12.0, 0.28, size=8, color=TEXT3, italic=True)

# ══════════════════════════════════════════════════════════
# SLIDE 11 — FORMATAÇÃO DE NÚMEROS
# ══════════════════════════════════════════════════════════
s = add_slide()
bg(s)
orange_bar(s)
section_label(s, "FORMATAÇÃO", 0.45, 0.2)
txt(s, "Formatação de números e percentuais", 0.45, 0.55, 10, 0.35, size=18, bold=True, color=TEXT)

r = rect(s, 0.45, 1.1, 12.4, 3.5, fill=CARD, line_color=BORDER, line_width=Pt(0.75))
txt(s, "Funções de formatação (padrão visao-financeira + painel-km)", 0.6, 1.18, 12.0, 0.3, size=11, bold=True, color=ORANGE)

code_fmt = (
    "const numFmt = v => {\n"
    "  const a = Math.abs(v), s = v < 0 ? '-' : '';\n"
    "  if (a >= 1e9) return s + (a/1e9).toFixed(2) + ' bi';  // ex: 1.23 bi\n"
    "  if (a >= 1e6) return s + (a/1e6).toFixed(2) + ' mi';  // ex: 5.70 mi\n"
    "  if (a >= 1e5) return s + (a/1e3).toFixed(2) + 'k';    // ex: 320.50k\n"
    "  return s + Math.round(a).toLocaleString('pt-BR');      // ex: 82.450\n"
    "};\n\n"
    "const fmt    = v => numFmt(v);\n"
    "const pctInt = v => (v >= 0 ? '+' : '') + v.toFixed(1) + '%';   // +11.4% ou -3.2%\n"
    "const pctR   = (a, b) => b ? ((a/b - 1)*100).toFixed(1) + '%' : '—';  // ratio %\n"
    "const pp     = v => v.toFixed(2) + ' pp';                        // pontos percentuais"
)
txt(s, code_fmt, 0.6, 1.55, 12.0, 2.9, size=9, color=TEXT2)

# Examples table
r = rect(s, 0.45, 4.8, 12.4, 2.35, fill=CARD, line_color=BORDER, line_width=Pt(0.75))
txt(s, "Exemplos de saída", 0.6, 4.88, 12.0, 0.3, size=11, bold=True, color=ORANGE)

examples_fmt = [
    ("1.234.567.890", "numFmt(v)", "1.23 bi"),
    ("5.700.000",     "numFmt(v)", "5.70 mi"),
    ("320.500",       "numFmt(v)", "320.50k"),
    ("82.450",        "numFmt(v)", "82.450"),
    ("11.4",          "pctInt(v)", "+11.4%"),
    ("-3.2",          "pctInt(v)", "-3.2%"),
    ("(905000, 850000)", "pctR(a,b)", "+6.5%"),
    ("2.35",          "pp(v)",    "2.35 pp"),
]
for i, (input_v, fn, output_v) in enumerate(examples_fmt):
    col = i % 4
    row = i // 4
    cx2 = 0.55 + col * 3.0
    cy2 = 5.22 + row * 0.5
    rb = rect(s, cx2, cy2, 2.85, 0.45, fill=CARD2 if (i // 4) % 2 else CARD)
    txt(s, fn, cx2+0.1, cy2+0.07, 1.3, 0.3, size=8, color=BLUE)
    txt(s, "→", cx2+1.4, cy2+0.08, 0.2, 0.28, size=9, color=TEXT3)
    txt(s, output_v, cx2+1.65, cy2+0.07, 1.1, 0.3, size=9, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)

# ══════════════════════════════════════════════════════════
# SLIDE 12 — MODO CLARO/ESCURO + RESUMO FINAL
# ══════════════════════════════════════════════════════════
s = add_slide()
bg(s)
orange_bar(s)
section_label(s, "MODO CLARO + RESUMO", 0.45, 0.2)
txt(s, "Tema + próximos passos", 0.45, 0.55, 10, 0.35, size=18, bold=True, color=TEXT)

# Theme section
r = rect(s, 0.45, 1.1, 6.0, 3.3, fill=CARD, line_color=BORDER, line_width=Pt(0.75))
txt(s, "Modo Claro / Escuro", 0.6, 1.18, 5.7, 0.3, size=11, bold=True, color=ORANGE)
theme_info = [
    ("Botão",    "lua/sol no header-right (theme-btn)"),
    ("Estado",   "localStorage 'bi_theme' = 'dark' | 'light'"),
    ("Header",   "SEMPRE escuro (não muda no modo claro)"),
    ("Body",     "body.light-mode .main { background: #F0F0F0 }"),
    ("Cards",    "body.light-mode .kpi-card { background: #FFF }"),
    ("Gráficos", "Regerar ao trocar tema — renderCharts(lastF)"),
    ("Ticks",    "getTick() retorna cor dinâmica por tema"),
]
for i, (k, v) in enumerate(theme_info):
    cy = 1.58 + i * 0.39
    txt(s, k + ":", 0.62, cy, 1.2, 0.32, size=9, bold=True, color=TEXT2)
    txt(s, v, 1.9, cy, 4.4, 0.32, size=9, color=TEXT)

# Pending items
r = rect(s, 6.7, 1.1, 6.2, 3.3, fill=CARD, line_color=RED, line_width=Pt(1))
txt(s, "⚠ Pendentes no main (push via MCP)", 6.85, 1.18, 5.9, 0.3, size=11, bold=True, color=RED)
pending = [
    ("index.html",                           "fix isApproved (crash hub) + pendente push"),
    ("visao-financeira/index.html",          "back-btn pendente push"),
    ("painel-km/index.html",                 "back-btn + números 2dec + pct 1dec — pendente push"),
    ("rs-por-km/index.html",                 "back-btn + font 13px — pendente push"),
    ("combustivel/arvore-combustivel/",      "back-btn — pendente push"),
]
for i, (file, status) in enumerate(pending):
    cy = 1.65 + i * 0.55
    txt(s, "• " + file, 6.85, cy, 5.9, 0.27, size=9, bold=True, color=ORANGE)
    txt(s, status, 6.85, cy+0.27, 5.9, 0.22, size=8, color=TEXT2)

# Roadmap
r = rect(s, 0.45, 4.6, 12.4, 2.7, fill=CARD, line_color=BORDER, line_width=Pt(0.75))
txt(s, "Roadmap de painéis a criar", 0.6, 4.68, 12.0, 0.3, size=11, bold=True, color=ORANGE)
todo = [
    ("Eficiência Km/L",         "/combustivel/eficiencia-kml/"),
    ("Preço R$/L",               "/combustivel/preco-litro/"),
    ("Consumo CO²",              "/combustivel/consumo-co2/"),
    ("Sub-hub Combustível",      "/combustivel/  (quando 2+ prontos)"),
    ("Ativação de Frota",        "/eficiencia-ativacao/"),
    ("Disponibilidade",          "/disponibilidade/"),
    ("Gerot / Auditorias / FCA", "/auditorias/  /fca/"),
    ("Prog. Reconhecimento",     "pasta a criar"),
    ("Aderência ao FCA",         "pasta a criar"),
    ("Painel de Metas",          "/painel-metas/"),
]
for i, (name, path) in enumerate(todo):
    col = i % 5
    row = i // 5
    cx = 0.55 + col * 2.4
    cy = 5.1 + row * 0.62
    rb = rect(s, cx, cy, 2.3, 0.55, fill=CARD2, line_color=BORDER, line_width=Pt(0.5))
    txt(s, name, cx+0.1, cy+0.04, 2.1, 0.27, size=8, bold=True, color=TEXT)
    txt(s, path, cx+0.1, cy+0.3,  2.1, 0.22, size=7, color=TEXT3, italic=True)

# ─── Save ──────────────────────────────────────────────────────────
out = '/home/user/Projeto-BI-App/PROJETO-BI-APRESENTACAO.pptx'
prs.save(out)
print(f"PPT salvo em: {out}")
print(f"Total de slides: {len(prs.slides)}")
